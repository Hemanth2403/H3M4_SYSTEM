from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "H3M4 Ecosystem: Collaborative Cybersecurity Framework for FinTech"
    subtitle.text = (
        "Project Presentation - Mid Semester Review\n\n"
        "Submitted to:\n"
        "School of Cyber Security & Digital Forensics,\n"
        "National Forensic Sciences University\n\n"
        "Name: M Hemanth naik\n"
        "Enrollment: 012200300004042\n"
        "Course: BTECH MTECH CS(CYBERSECURITY)\n\n"
        "Guide: Dr. Ramya Shah, Assistant Professor"
    )

    # --- Slide 2: Table of Contents ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "• Abstract"
    p = tf.add_paragraph()
    p.text = "• Literature Review (LR)"
    p = tf.add_paragraph()
    p.text = "• Motivation"
    p = tf.add_paragraph()
    p.text = "• Problem Statement"
    p = tf.add_paragraph()
    p.text = "• Project Objectives"
    p = tf.add_paragraph()
    p.text = "• Approach & Methodology"
    p = tf.add_paragraph()
    p.text = "• Implementation & Tools"
    p = tf.add_paragraph()
    p.text = "• Project Architecture (Graphical)"

    # --- Slide 3: Abstract ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Abstract"
    tf = slide.placeholders[1].text_frame
    # Abstract Slide often needs justification
    tf.text = "H3M4 is a collaborative cybersecurity framework designed to bridge the trust gap between Security Researchers, Financial Enterprises, and Law Enforcement."
    tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
    p = tf.add_paragraph()
    p.text = "Key features include a live 'Global Signal Graph' for threat visualization, a specialized 'Breach Bot' for PII leak discovery, and an automated Intelligence-to-Enforcement pipeline."
    p.alignment = PP_ALIGN.JUSTIFY
    p = tf.add_paragraph()
    p.text = "The system ensures immutable evidence handling and sector-wide 'herd immunity' against sophisticated cyber threats."
    p.alignment = PP_ALIGN.JUSTIFY

    # --- Slide 4: Literature Review - Findings & Limitations (Part 1) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Literature Review: Findings"
    tf = slide.placeholders[1].text_frame
    def add_justified(text_frame, text):
        p = text_frame.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.JUSTIFY
        return p

    add_justified(tf, "• AI-Driven Threats (2023): Research shows 11 central threats requiring 24/7 SIEM monitoring.")
    add_justified(tf, "• Rising Incident Volume (2024): 53% surge in FinTech attacks; AI/ML is crucial for real-time detection.")
    add_justified(tf, "• Strategic Intelligence (2025): 91% of security leaders are moving towards intelligence-guided investments.")
    add_justified(tf, "• LLMs in SecOps: Large Language Models are automating the parsing of dark web forum data.")

    # --- Slide 5: Literature Review - Limitations ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Literature Review: Limitations"
    tf = slide.placeholders[1].text_frame
    add_justified(tf, "• Black Box Defense: Traditional SIEMs lack cross-sector transparency, making it hard to stop widespread campaigns.")
    add_justified(tf, "• Operational Overhead: Moving to Zero-Trust or DORA-compliant models involves high technical debt.")
    add_justified(tf, "• Translation Silos: Technical indicators (IOCs) are often too complex for business stakeholders to act upon.")
    add_justified(tf, "• AI Hallucinations: Over-reliance on LLMs for hacker attribution requires human-in-the-loop validation.")

    # --- Slide 6: Motivation ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Motivation"
    tf = slide.placeholders[1].text_frame
    add_justified(tf, "• The 53% Incident Spike: Real-world reports of surging attacks in the financial sector.")
    add_justified(tf, "• The 'Hacker-to-Justice' Gap: No secure bridge for ethical hackers to report directly to Law Enforcement.")
    add_justified(tf, "• Data Silos: Need for a 'Federated Defense' where one bank's breach detection becomes every bank's protection.")
    add_justified(tf, "• Cyber-Forensics Efficiency: Reducing the time taken to generate court-admissible evidence from months to hours.")

    # --- Slide 7: Problem Statement ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement"
    tf = slide.placeholders[1].text_frame
    add_justified(tf, "Critical Project Question: 'How can we create a trust-based ecosystem where rivals (banks) and external actors (researchers) collaborate to defeat shared cyber threats?'")
    add_justified(tf, "Identified Limitations in Current Systems:")
    add_justified(tf, "  - High Reputation Risk for banks reporting breaches anonymously.")
    add_justified(tf, "  - Lack of integration between Researcher research and Police FIRs.")
    add_justified(tf, "  - Massive growth of PII leaks in Telegram/Dark Web without scalable traceability.")

    # --- Slide 8: Project Objectives ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Objectives"
    tf = slide.placeholders[1].text_frame
    add_justified(tf, "1. Develop Multi-Role Hub: Create specialized workspaces for Researchers (Sources), Enterprises (Sensors), and Police (Enforcers).")
    add_justified(tf, "2. Build Global Threat Graph: Implement real-time visualization of Actors, Vectors, and Targets using graph-based rendering.")
    add_justified(tf, "3. Automate Forensic Alignment: Enabling instant attachment of verified technical intelligence to legal case files (FIRs).")

    # --- Slide 9: Approach & Methodology ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Approach & Methodology"
    tf = slide.placeholders[1].text_frame
    add_justified(tf, "• Role-Based Access Control (RBAC): Ensuring strict boundary separation and data privacy.")
    add_justified(tf, "• Zero-Knowledge Principles: Allowing enterprises to share attack signatures without exposing user PII.")
    add_justified(tf, "• Reactive & Proactive Hybrid: Combining live SIEM log analysis (Reactive) with vulnerability research (Proactive).")
    add_justified(tf, "• Immutable Chain of Custody: Ensuring all evidence sharing is cryptographically logged.")

    # --- Slide 10: Implementation & Tools ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Implementation & Tools"
    tf = slide.placeholders[1].text_frame
    add_justified(tf, "• Frontend: React.js, TypeScript, Tailwind CSS, Framer Motion.")
    add_justified(tf, "• Backend: Node.js, Express, Passport.js, React Query.")
    add_justified(tf, "• Analysis: Wazuh/SIEM Integration, 'Breach Bot' Python wrappers.")
    add_justified(tf, "• Data Layer: Database simulation for high-performance demonstration.")
    add_justified(tf, "• Lab Setup: Multi-node development environment for cross-role simulation.")

    # --- Slide 11: Project Architecture (Graphical) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text_frame.text = "Project Architecture: Collaborative Ecosystem Cycle"
    
    # Simple diagram using shapes
    def add_node(left, top, text, color):
        shape = slide.shapes.add_shape(6, Inches(left), Inches(top), Inches(2), Inches(0.8)) # 6 is RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.text = text
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return shape

    # Define nodes
    researcher = add_node(4, 1.5, "Researcher\n(Intelligence Source)", RGBColor(0, 102, 204))
    admin = add_node(4, 3.0, "System Admin\n(Validator/Governor)", RGBColor(128, 0, 128))
    enterprise = add_node(1.5, 4.5, "Enterprise Node\n(Defender)", RGBColor(0, 153, 0))
    police = add_node(6.5, 4.5, "Police Command\n(Enforcer)", RGBColor(204, 0, 0))

    # Add arrows (very basic representation)
    # Researcher -> Admin
    prs.slides[10].shapes.add_connector(2, researcher.left + Inches(1), researcher.top + Inches(0.8), admin.left + Inches(1), admin.top) # STRAIGHT arrow
    
    # Admin -> Enterprise & Admin -> Police
    prs.slides[10].shapes.add_connector(2, admin.left + Inches(0.5), admin.top + Inches(0.8), enterprise.left + Inches(1), enterprise.top)
    prs.slides[10].shapes.add_connector(2, admin.left + Inches(1.5), admin.top + Inches(0.8), police.left + Inches(1), police.top)

    # Enterprise (Victim/Evidence) -> Police (FIR Case)
    prs.slides[10].shapes.add_connector(2, enterprise.left + Inches(2), enterprise.top + Inches(0.4), police.left, police.top + Inches(0.4))

    # --- Slide 12: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You!"
    slide.placeholders[1].text = "Questions & Suggestions are welcome."

    # Save the presentation
    prs.save("H3M4_Project_Presentation.pptx")
    print("Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
